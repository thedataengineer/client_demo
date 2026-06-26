import os
from docx import Document
from openpyxl import Workbook
from pptx import Presentation

os.makedirs('docs', exist_ok=True)

# 1. Create Word Document
doc = Document()
doc.add_heading('Product Requirements Document (PRD)', 0)

doc.add_heading('Overview', level=1)
doc.add_paragraph('This application is a full-stack dashboard utilizing React for the frontend, FastAPI for the backend, and PostgreSQL for the database.')

doc.add_heading('Features', level=1)
doc.add_paragraph('1. Real-time data visualization.')
doc.add_paragraph('2. RESTful API integration.')
doc.add_paragraph('3. Persistent storage with PostgreSQL.')

doc.save('docs/requirements.docx')
print("Generated requirements.docx")

# 2. Create Excel Spreadsheet
wb = Workbook()
ws = wb.active
ws.title = "API Endpoints"

ws.append(["Endpoint", "Method", "Description"])
ws.append(["/api/items", "GET", "Fetch all items from the database"])
ws.append(["/api/items", "POST", "Create a new item"])
ws.append(["/", "GET", "Health check"])

wb.save('docs/api_endpoints.xlsx')
print("Generated api_endpoints.xlsx")

# 3. Create PowerPoint Presentation
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Project Overview"
subtitle.text = "React + FastAPI + PostgreSQL"

bullet_slide_layout = prs.slide_layouts[1]
slide2 = prs.slides.add_slide(bullet_slide_layout)
shapes = slide2.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Architecture"
tf = body_shape.text_frame
tf.text = "Frontend: Vite + React"

p = tf.add_paragraph()
p.text = "Backend: Python + FastAPI"

p2 = tf.add_paragraph()
p2.text = "Database: PostgreSQL"

p3 = tf.add_paragraph()
p3.text = "Infra: Terraform + Docker"

prs.save('docs/presentation.pptx')
print("Generated presentation.pptx")
